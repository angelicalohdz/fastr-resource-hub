#!/usr/bin/env python3
"""
═══════════════════════════════════════════════════════════════════════
              FASTR POWERPOINT CONVERTER
═══════════════════════════════════════════════════════════════════════

Converts Marp markdown decks to branded PowerPoint presentations.

USAGE:
    python3 tools/03_convert_pptx.py                           # Interactive
    python3 tools/03_convert_pptx.py outputs/workshop_deck.md  # Direct

FEATURES:
  ✓ FASTR brand colors and styling
  ✓ Proper slide layouts (title, content, agenda, two-column)
  ✓ Tables with styled headers
  ✓ Images embedded correctly
  ✓ Break slides with centered text

═══════════════════════════════════════════════════════════════════════
"""

import argparse
import os
import re
import sys
from pathlib import Path

# ═══════════════════════════════════════════════════════════════════════════════
# AUTO-DETECT AND USE VENV
# ═══════════════════════════════════════════════════════════════════════════════

def ensure_venv():
    """Re-execute with venv Python if not already in venv."""
    if sys.prefix != sys.base_prefix:
        return
    script_dir = Path(__file__).resolve().parent
    project_root = script_dir.parent
    for venv_name in ['.venv', 'venv']:
        venv_python = project_root / venv_name / 'bin' / 'python3'
        if venv_python.exists():
            os.execv(str(venv_python), [str(venv_python)] + sys.argv)

ensure_venv()

# Now import python-pptx (after venv activation)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx not installed.")
    print("Run: pip install python-pptx Pillow")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    Image = None

# ═══════════════════════════════════════════════════════════════════════════════
# FASTR BRAND CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════

class Colors:
    """FASTR brand colors from fastr-theme.css"""
    DEEP_GREEN = RGBColor(0x09, 0x54, 0x4F)   # #09544F - H1
    DARK_GREEN = RGBColor(0x0C, 0x71, 0x6B)   # #0C716B
    GREEN = RGBColor(0x1F, 0x9A, 0x9C)        # #1F9A9C
    LIME = RGBColor(0xD0, 0xCB, 0x17)         # #D0CB17 - H1 underline
    NAVY = RGBColor(0x21, 0x56, 0x8C)         # #21568C - H2
    BLUE = RGBColor(0x1A, 0x90, 0xC0)         # #1A90C0 - H2 underline
    LIGHT_BLUE = RGBColor(0xCA, 0xE6, 0xE9)   # #CAE6E9 - table headers

    GOLD = RGBColor(0xD8, 0xA8, 0x22)         # #D8A822
    PURPLE = RGBColor(0x7A, 0x1F, 0x6E)       # #7A1F6E
    ORCHID = RGBColor(0xBD, 0x50, 0x91)       # #BD5091
    CORAL = RGBColor(0xFF, 0x64, 0x62)        # #FF6462

    TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)    # #2c3e50
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)    # #333333 - body text
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x00, 0x00, 0x00)


class Fonts:
    """Font settings"""
    FAMILY = 'Arial'  # Universal fallback
    H1_SIZE = Pt(40)
    H2_SIZE = Pt(32)
    H3_SIZE = Pt(24)
    BODY_SIZE = Pt(18)
    TABLE_SIZE = Pt(14)
    SMALL_SIZE = Pt(12)


class Layout:
    """Slide dimensions and margins (16:9)"""
    WIDTH = Inches(13.333)
    HEIGHT = Inches(7.5)
    MARGIN_LEFT = Inches(0.5)
    MARGIN_RIGHT = Inches(0.5)
    MARGIN_TOP = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)
    CONTENT_WIDTH = Inches(12.333)


# ═══════════════════════════════════════════════════════════════════════════════
# MARKDOWN PARSER
# ═══════════════════════════════════════════════════════════════════════════════

def parse_markdown(content):
    """
    Parse Marp markdown into slide data structures.

    Returns list of dicts with:
    - raw: original markdown
    - headers: [(level, text), ...]
    - bullets: [text, ...]
    - table: [[row], ...] or None
    - images: [{'alt': str, 'path': str}, ...]
    - columns: {'left': str, 'right': str} or None
    - css_class: from <!-- _class: xxx -->
    """
    # Strip YAML frontmatter
    if content.startswith('---'):
        match = re.match(r'^---\n.*?\n---\n?', content, re.DOTALL)
        if match:
            content = content[match.end():]

    # Split into slides
    raw_slides = re.split(r'\n---\s*\n', content)

    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        slide = {
            'raw': raw,
            'headers': [],
            'bullets': [],
            'paragraphs': [],  # Plain text lines (not headers, bullets, or tables)
            'content': [],     # Ordered list of ('paragraph', text) or ('bullet', text)
            'table': None,
            'images': [],
            'columns': None,
            'css_class': None,
            'html_content': None,
        }

        # Extract CSS class directive
        class_match = re.search(r'<!--\s*_class:\s*(\w+)\s*-->', raw)
        if class_match:
            slide['css_class'] = class_match.group(1)
            raw = re.sub(r'<!--\s*_class:\s*\w+\s*-->', '', raw)

        # Extract headers
        for match in re.finditer(r'^(#{1,6})\s+(.+)$', raw, re.MULTILINE):
            level = len(match.group(1))
            text = match.group(2).strip()
            slide['headers'].append((level, text))

        # Extract images ![alt](path)
        for match in re.finditer(r'!\[([^\]]*)\]\(([^)]+)\)', raw):
            slide['images'].append({
                'alt': match.group(1),
                'path': match.group(2).split()[0]  # Remove any title
            })

        # Extract table
        table_lines = []
        in_table = False
        for line in raw.split('\n'):
            if '|' in line and line.strip().startswith('|'):
                in_table = True
                # Skip separator line
                if re.match(r'^\|[\s\-:|]+\|$', line.strip()):
                    continue
                cells = [c.strip() for c in line.strip().strip('|').split('|')]
                table_lines.append(cells)
            elif in_table:
                break
        if table_lines:
            slide['table'] = table_lines

        # Extract columns div
        cols_match = re.search(
            r'<div\s+class="columns[^"]*">\s*<div>(.*?)</div>\s*<div>(.*?)</div>\s*</div>',
            raw, re.DOTALL
        )
        if cols_match:
            slide['columns'] = {
                'left': cols_match.group(1).strip(),
                'right': cols_match.group(2).strip()
            }

        # Check for other HTML content (like styled img tags)
        if '<img' in raw:
            slide['html_content'] = raw

        # Extract content IN ORDER - headers (H3+), bullets, numbered lists, and paragraphs
        for line in raw.split('\n'):
            line_stripped = line.strip()
            if not line_stripped:
                continue
            # Include H3+ headers in content flow (H1/H2 rendered separately at top)
            header_match = re.match(r'^(#{3,6})\s+(.+)$', line_stripped)
            if header_match:
                level = len(header_match.group(1))
                text = header_match.group(2).strip()
                slide['content'].append(('header', (level, text)))
                continue
            # Skip H1/H2 headers (rendered at top)
            if line_stripped.startswith('#'):
                continue
            # Skip table lines
            if line_stripped.startswith('|'):
                continue
            # Skip HTML
            if line_stripped.startswith('<') or line_stripped.startswith('!'):
                continue
            # Skip comments
            if line_stripped.startswith('<!--'):
                continue

            # Check for bullet
            bullet_match = re.match(r'^[-*]\s+(.+)$', line_stripped)
            if bullet_match:
                text = bullet_match.group(1).strip()
                slide['bullets'].append(text)
                slide['content'].append(('bullet', text))
                continue

            # Check for numbered list
            num_match = re.match(r'^\d+\.\s+(.+)$', line_stripped)
            if num_match:
                text = num_match.group(1).strip()
                slide['bullets'].append(text)
                slide['content'].append(('bullet', text))
                continue

            # It's a paragraph
            slide['paragraphs'].append(line_stripped)
            slide['content'].append(('paragraph', line_stripped))

        slides.append(slide)

    return slides


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE TYPE DETECTOR
# ═══════════════════════════════════════════════════════════════════════════════

def detect_slide_type(slide, index):
    """
    Detect slide type from content.

    Returns: 'title', 'agenda', 'break', 'two_column', 'table', 'image', 'content'
    """
    headers = slide['headers']
    h1_text = headers[0][1] if headers and headers[0][0] == 1 else ''

    # Title slide: first slide with logo
    if index == 0:
        for img in slide['images']:
            if 'logo' in img['path'].lower() or 'fastr' in img['path'].lower():
                return 'title'

    # Agenda slide
    if slide['css_class'] == 'agenda' or (slide['table'] and 'agenda' in h1_text.lower()):
        return 'agenda'

    # Break slide: emoji or "break" in H1
    break_emojis = ['☕', '🍽', '🌙', '🎉', '👋', '⏰']
    if any(emoji in h1_text for emoji in break_emojis):
        return 'break'
    if re.search(r'\b(break|lunch|tea)\b', h1_text, re.IGNORECASE):
        return 'break'

    # Two-column slide
    if slide['columns']:
        return 'two_column'

    # Table slide (non-agenda)
    if slide['table']:
        return 'table'

    # Image-heavy slide
    if slide['images'] and len(slide['bullets']) < 2:
        return 'image'

    # Default: content slide
    return 'content'


# ═══════════════════════════════════════════════════════════════════════════════
# IMAGE HANDLING
# ═══════════════════════════════════════════════════════════════════════════════

def resolve_image_path(img_path, base_dir, md_dir):
    """
    Resolve relative image path to absolute path.
    Tries multiple locations.
    """
    if img_path.startswith(('http://', 'https://')):
        return None  # Skip URLs

    # Clean up the path
    img_path = img_path.replace('%20', ' ')

    paths_to_try = [
        os.path.join(md_dir, img_path),
        os.path.join(base_dir, img_path.lstrip('../')),
        os.path.join(base_dir, 'resources', 'logos', os.path.basename(img_path)),
        os.path.join(base_dir, 'resources', 'diagrams', os.path.basename(img_path)),
        os.path.join(base_dir, 'resources', 'default_outputs', os.path.basename(img_path)),
        os.path.join(base_dir, 'assets', os.path.basename(img_path)),
    ]

    # Handle ../resources/... paths from outputs/ folder
    if '../resources/' in img_path:
        clean_path = img_path.replace('../', '')
        paths_to_try.insert(0, os.path.join(base_dir, clean_path))

    for path in paths_to_try:
        if os.path.exists(path):
            return os.path.abspath(path)

    return None


def convert_svg_to_png(svg_path, temp_dir):
    """Convert SVG to PNG for PowerPoint compatibility."""
    try:
        import cairosvg
        png_path = os.path.join(temp_dir, Path(svg_path).stem + '.png')
        cairosvg.svg2png(url=svg_path, write_to=png_path, scale=2.0)
        return png_path
    except ImportError:
        return None
    except Exception:
        return None


def add_image_to_slide(slide, img_path, left, top, width=None, height=None):
    """Add image to slide, handling SVG conversion if needed."""
    if not img_path or not os.path.exists(img_path):
        return None

    # Convert SVG if needed
    if img_path.lower().endswith('.svg'):
        import tempfile
        with tempfile.TemporaryDirectory() as temp_dir:
            png_path = convert_svg_to_png(img_path, temp_dir)
            if png_path:
                img_path = png_path
            else:
                return None  # SVG conversion failed

    try:
        if width and height:
            return slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            return slide.shapes.add_picture(img_path, left, top)
    except Exception as e:
        print(f"   Warning: Could not add image: {e}")
        return None


# ═══════════════════════════════════════════════════════════════════════════════
# TEXT STYLING HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def style_paragraph(para, font_size, color, bold=False, italic=False):
    """Apply styling to a paragraph."""
    para.font.size = font_size
    para.font.color.rgb = color
    para.font.bold = bold
    para.font.italic = italic
    para.font.name = Fonts.FAMILY


def add_text_box(slide, left, top, width, height, text, font_size, color,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a styled text box to slide."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Set vertical alignment
    try:
        tf.anchor = anchor
    except:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    style_paragraph(p, font_size, color, bold)

    return shape


def add_h1_with_underline(slide, text, top=None):
    """Add H1 header with lime underline."""
    if top is None:
        top = Layout.MARGIN_TOP

    # Add title text
    title_shape = add_text_box(
        slide,
        Layout.MARGIN_LEFT, top,
        Layout.CONTENT_WIDTH, Inches(0.8),
        text, Fonts.H1_SIZE, Colors.DEEP_GREEN, bold=True
    )

    # Add underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Layout.MARGIN_LEFT, top + Inches(0.7),
        Inches(6), Inches(0.06)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = Colors.LIME
    underline.line.fill.background()

    return title_shape


def add_h2_with_underline(slide, text, top):
    """Add H2 header with blue underline."""
    title_shape = add_text_box(
        slide,
        Layout.MARGIN_LEFT, top,
        Layout.CONTENT_WIDTH, Inches(0.6),
        text, Fonts.H2_SIZE, Colors.NAVY, bold=True
    )

    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Layout.MARGIN_LEFT, top + Inches(0.55),
        Inches(5), Inches(0.04)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = Colors.BLUE
    underline.line.fill.background()

    return title_shape


def add_bullet_list(slide, bullets, left, top, width, font_size=None):
    """Add a bullet list to slide."""
    if not bullets:
        return None

    if font_size is None:
        font_size = Fonts.BODY_SIZE

    height = Inches(0.4 * len(bullets) + 0.2)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Clean up bullet text (remove markdown formatting)
        text = re.sub(r'\*\*([^*]+)\*\*', r'\1', bullet)  # Bold
        text = re.sub(r'\*([^*]+)\*', r'\1', text)  # Italic
        text = re.sub(r'`([^`]+)`', r'\1', text)  # Code

        p.text = f"• {text}"
        p.level = 0
        style_paragraph(p, font_size, Colors.TEXT_DARK)

    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_title_slide(prs, data, base_dir, md_dir):
    """Build title slide with centered content and logo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Get title
    title = data['headers'][0][1] if data['headers'] else 'FASTR Workshop'

    # Centered title
    add_text_box(
        slide,
        Inches(1), Inches(2),
        Inches(11.333), Inches(1.2),
        title, Fonts.H1_SIZE, Colors.DEEP_GREEN,
        bold=True, align=PP_ALIGN.CENTER
    )

    # Underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3), Inches(3.2),
        Inches(7.333), Inches(0.06)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = Colors.LIME
    underline.line.fill.background()

    # Extract subtitle info from raw content
    raw = data['raw']

    # Look for date/location line (bold text)
    date_match = re.search(r'\*\*([^*]+)\*\*\s*\|\s*\*\*([^*]+)\*\*', raw)
    if date_match:
        subtitle = f"{date_match.group(1)} | {date_match.group(2)}"
        add_text_box(
            slide,
            Inches(1), Inches(3.5),
            Inches(11.333), Inches(0.5),
            subtitle, Fonts.H2_SIZE, Colors.NAVY,
            align=PP_ALIGN.CENTER
        )

    # Look for facilitator (italic text)
    fac_match = re.search(r'\*([^*]+)\*(?!\*)', raw)
    if fac_match and 'Facilitator' in fac_match.group(1) or fac_match:
        add_text_box(
            slide,
            Inches(1), Inches(4.2),
            Inches(11.333), Inches(0.4),
            fac_match.group(1) if fac_match else '', Fonts.BODY_SIZE, Colors.TEXT_DARK,
            align=PP_ALIGN.CENTER
        )

    # Add logo
    for img in data['images']:
        img_path = resolve_image_path(img['path'], base_dir, md_dir)
        if img_path and ('logo' in img_path.lower() or 'fastr' in img_path.lower()):
            add_image_to_slide(
                slide, img_path,
                Inches(10.5), Inches(6.2),
                width=Inches(2.3)
            )
            break

    return slide


def build_agenda_slide(prs, data, base_dir, md_dir):
    """Build agenda slide with compact table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    title = data['headers'][0][1] if data['headers'] else 'Workshop Agenda'
    add_h1_with_underline(slide, title)

    # Table
    if data['table']:
        table_data = data['table']
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 3

        table = slide.shapes.add_table(
            rows, cols,
            Layout.MARGIN_LEFT, Inches(1.3),
            Inches(12), Inches(0.35 * rows)
        ).table

        # Set column widths
        if cols >= 3:
            table.columns[0].width = Inches(2.5)   # Time
            table.columns[1].width = Inches(7)     # Session
            table.columns[2].width = Inches(2.5)   # Speaker

        # Fill table
        for r_idx, row in enumerate(table_data):
            is_header = r_idx == 0
            for c_idx, cell_text in enumerate(row):
                cell = table.cell(r_idx, c_idx)

                # Clean text
                text = re.sub(r'\*\*([^*]+)\*\*', r'\1', cell_text)
                text = re.sub(r'\*([^*]+)\*', r'\1', text)

                cell.text = text

                # Style cell
                para = cell.text_frame.paragraphs[0]
                para.font.size = Fonts.TABLE_SIZE
                para.font.name = Fonts.FAMILY

                if is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = Colors.LIGHT_BLUE
                    para.font.color.rgb = Colors.NAVY
                    para.font.bold = True
                else:
                    para.font.color.rgb = Colors.TEXT_DARK

    return slide


def build_break_slide(prs, data, base_dir, md_dir):
    """Build break slide with large centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Get title (with emoji)
    title = data['headers'][0][1] if data['headers'] else 'Break'

    # Large centered title
    add_text_box(
        slide,
        Inches(1), Inches(2.8),
        Inches(11.333), Inches(1.5),
        title, Pt(56), Colors.DEEP_GREEN,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )

    # Look for time info in raw content
    raw = data['raw']
    time_match = re.search(r'(\d+)\s*minutes?', raw, re.IGNORECASE)
    resume_match = re.search(r'(?:resume|back|return)[^\d]*(\d+[:\d]*\s*(?:AM|PM)?)', raw, re.IGNORECASE)

    subtitle_parts = []
    if time_match:
        subtitle_parts.append(f"{time_match.group(1)} minutes")
    if resume_match:
        subtitle_parts.append(f"Back at {resume_match.group(1)}")

    if subtitle_parts:
        add_text_box(
            slide,
            Inches(1), Inches(4.3),
            Inches(11.333), Inches(0.6),
            " • ".join(subtitle_parts), Fonts.H2_SIZE, Colors.NAVY,
            align=PP_ALIGN.CENTER
        )

    return slide


def build_two_column_slide(prs, data, base_dir, md_dir):
    """Build two-column slide with text and image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    if data['headers']:
        title = data['headers'][0][1]
        if data['headers'][0][0] == 1:
            add_h1_with_underline(slide, title)
        else:
            add_h2_with_underline(slide, title, Layout.MARGIN_TOP)

    # Parse column content
    left_content = data['columns']['left']
    right_content = data['columns']['right']

    # Determine which side has image
    left_has_image = '![' in left_content or '<img' in left_content
    right_has_image = '![' in right_content or '<img' in right_content

    content_top = Inches(1.4)
    col_width = Inches(5.8)
    col_height = Inches(5.5)

    # Left column
    if left_has_image:
        # Add image
        img_match = re.search(r'!\[[^\]]*\]\(([^)]+)\)', left_content)
        if img_match:
            img_path = resolve_image_path(img_match.group(1).split()[0], base_dir, md_dir)
            if img_path:
                add_image_to_slide(
                    slide, img_path,
                    Layout.MARGIN_LEFT, content_top,
                    width=col_width
                )
    else:
        # Add bullets
        bullets = re.findall(r'^[-*]\s+(.+)$', left_content, re.MULTILINE)
        if bullets:
            add_bullet_list(slide, bullets, Layout.MARGIN_LEFT, content_top, col_width)

    # Right column
    right_left = Inches(7)
    if right_has_image:
        img_match = re.search(r'!\[[^\]]*\]\(([^)]+)\)', right_content)
        if img_match:
            img_path = resolve_image_path(img_match.group(1).split()[0], base_dir, md_dir)
            if img_path:
                add_image_to_slide(
                    slide, img_path,
                    right_left, content_top,
                    width=col_width
                )
    else:
        bullets = re.findall(r'^[-*]\s+(.+)$', right_content, re.MULTILINE)
        if bullets:
            add_bullet_list(slide, bullets, right_left, content_top, col_width)

    return slide


def build_table_slide(prs, data, base_dir, md_dir):
    """Build slide with table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    if data['headers']:
        title = data['headers'][0][1]
        if data['headers'][0][0] == 1:
            add_h1_with_underline(slide, title)
        else:
            add_h2_with_underline(slide, title, Layout.MARGIN_TOP)

    # Table
    if data['table']:
        table_data = data['table']
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 2

        # Calculate table dimensions
        table_width = min(Inches(12), Inches(2.5 * cols))
        row_height = Inches(0.4)

        table = slide.shapes.add_table(
            rows, cols,
            Layout.MARGIN_LEFT, Inches(1.4),
            table_width, row_height * rows
        ).table

        # Fill table
        for r_idx, row in enumerate(table_data):
            is_header = r_idx == 0
            for c_idx, cell_text in enumerate(row):
                if c_idx >= cols:
                    continue
                cell = table.cell(r_idx, c_idx)

                # Clean text
                text = re.sub(r'\*\*([^*]+)\*\*', r'\1', cell_text)
                text = re.sub(r'\*([^*]+)\*', r'\1', text)

                cell.text = text

                para = cell.text_frame.paragraphs[0]
                para.font.size = Fonts.BODY_SIZE
                para.font.name = Fonts.FAMILY

                if is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = Colors.LIGHT_BLUE
                    para.font.color.rgb = Colors.NAVY
                    para.font.bold = True
                else:
                    para.font.color.rgb = Colors.TEXT_DARK

    return slide


def build_image_slide(prs, data, base_dir, md_dir):
    """Build image-focused slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    if data['headers']:
        title = data['headers'][0][1]
        if data['headers'][0][0] == 1:
            add_h1_with_underline(slide, title)
        else:
            add_h2_with_underline(slide, title, Layout.MARGIN_TOP)

    # Add main image
    if data['images']:
        img = data['images'][0]
        img_path = resolve_image_path(img['path'], base_dir, md_dir)
        if img_path:
            # Center the image
            add_image_to_slide(
                slide, img_path,
                Inches(1.5), Inches(1.5),
                width=Inches(10)
            )

    # Add any bullets below
    if data['bullets']:
        add_bullet_list(
            slide, data['bullets'],
            Layout.MARGIN_LEFT, Inches(5.5),
            Layout.CONTENT_WIDTH, Fonts.BODY_SIZE
        )

    return slide


def build_content_slide(prs, data, base_dir, md_dir):
    """Build standard content slide with headers and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    current_top = Layout.MARGIN_TOP

    # Add H1/H2 headers at top (H3+ are in content flow)
    for level, text in data['headers']:
        if level == 1:
            add_h1_with_underline(slide, text, current_top)
            current_top += Inches(0.9)
        elif level == 2:
            add_h2_with_underline(slide, text, current_top)
            current_top += Inches(0.75)
        # H3+ handled in content loop below

    # Render ALL content (H3+, paragraphs, bullets) in ONE text box
    content = data.get('content', [])
    if content:
        # Calculate height based on content
        height = Inches(0.35 * len(content) + 0.5)
        shape = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, current_top + Inches(0.1),
            Layout.CONTENT_WIDTH, height
        )
        tf = shape.text_frame
        tf.word_wrap = True

        first_para = True
        for item_type, text in content:
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()

            if item_type == 'header':
                level, header_text = text
                p.text = header_text
                style_paragraph(p, Fonts.H3_SIZE, Colors.PURPLE, bold=True)

            elif item_type == 'bullet':
                # Clean markdown formatting
                clean_text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
                clean_text = re.sub(r'\*([^*]+)\*', r'\1', clean_text)
                p.text = f"• {clean_text}"
                style_paragraph(p, Fonts.BODY_SIZE, Colors.TEXT_DARK)

            elif item_type == 'paragraph':
                display_text = text
                is_bold = False
                if text.startswith('**') and text.endswith('**'):
                    display_text = text[2:-2]
                    is_bold = True
                p.text = display_text
                style_paragraph(p, Fonts.BODY_SIZE, Colors.DARK_GRAY, bold=is_bold)

    # Add image if present (to the right)
    if data['images']:
        img = data['images'][0]
        img_path = resolve_image_path(img['path'], base_dir, md_dir)
        if img_path:
            add_image_to_slide(
                slide, img_path,
                Inches(8), Inches(1.5),
                width=Inches(4.5)
            )

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN CONVERTER
# ═══════════════════════════════════════════════════════════════════════════════

def convert_to_pptx(md_file, base_dir, output_path=None):
    """
    Convert Marp markdown to PowerPoint.

    Returns True on success, False on failure.
    """
    # Validate input
    if not os.path.exists(md_file):
        print(f"Error: File not found: {md_file}")
        return False

    print("\n" + "=" * 70)
    print("           CONVERTING TO POWERPOINT")
    print("=" * 70)

    # Read markdown
    print(f"\n   Reading: {os.path.basename(md_file)}")
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Parse markdown
    print("   Parsing slides...")
    slides_data = parse_markdown(content)
    print(f"   Found {len(slides_data)} slides")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Layout.WIDTH
    prs.slide_height = Layout.HEIGHT

    # Get paths
    md_dir = os.path.dirname(os.path.abspath(md_file))

    # Build each slide
    print("   Building slides...")
    builders = {
        'title': build_title_slide,
        'agenda': build_agenda_slide,
        'break': build_break_slide,
        'two_column': build_two_column_slide,
        'table': build_table_slide,
        'image': build_image_slide,
        'content': build_content_slide,
    }

    type_counts = {}
    for i, data in enumerate(slides_data):
        slide_type = detect_slide_type(data, i)
        type_counts[slide_type] = type_counts.get(slide_type, 0) + 1

        builder = builders.get(slide_type, build_content_slide)
        try:
            builder(prs, data, base_dir, md_dir)
        except Exception as e:
            print(f"   Warning: Error building slide {i+1}: {e}")
            # Build as content slide fallback
            try:
                build_content_slide(prs, data, base_dir, md_dir)
            except:
                pass

    # Show type breakdown
    print("   Slide types:")
    for stype, count in sorted(type_counts.items()):
        print(f"      {stype}: {count}")

    # Save
    if output_path is None:
        output_path = md_file.replace('.md', '.pptx')

    print(f"\n   Saving: {os.path.basename(output_path)}")
    prs.save(output_path)

    file_size = os.path.getsize(output_path) / 1024

    print("\n" + "=" * 70)
    print("                    SUCCESS!")
    print("=" * 70)
    print(f"\n   Output: {output_path}")
    print(f"   Size: {file_size:.1f} KB")
    print(f"   Slides: {len(prs.slides)}")
    print("\n" + "=" * 70 + "\n")

    return True


# ═══════════════════════════════════════════════════════════════════════════════
# CLI INTERFACE
# ═══════════════════════════════════════════════════════════════════════════════

def list_available_decks(base_dir):
    """List markdown decks in outputs/ folder."""
    outputs_dir = os.path.join(base_dir, "outputs")
    if not os.path.exists(outputs_dir):
        return []

    decks = []
    for file in os.listdir(outputs_dir):
        if file.endswith('.md') and not file.startswith('.'):
            decks.append(file)

    return sorted(decks)


def prompt_for_deck(base_dir):
    """Interactive mode: ask user which deck to convert."""
    print("\n" + "=" * 70)
    print("              AVAILABLE DECKS")
    print("=" * 70 + "\n")

    decks = list_available_decks(base_dir)

    if not decks:
        print("No decks found in outputs/ folder!")
        print("\nBuild a deck first:")
        print("   python3 tools/02_build_deck.py --workshop YOUR-WORKSHOP")
        sys.exit(1)

    for i, deck in enumerate(decks, 1):
        deck_path = os.path.join(base_dir, "outputs", deck)
        size = os.path.getsize(deck_path) / 1024
        print(f"  {i}. {deck} ({size:.1f} KB)")

    print("\n" + "-" * 70)

    while True:
        try:
            choice = input("\nWhich deck to convert? (number or name): ").strip()

            if choice.isdigit():
                idx = int(choice) - 1
                if 0 <= idx < len(decks):
                    return decks[idx]

            if choice in decks:
                return choice

            if not choice.endswith('.md'):
                if choice + '.md' in decks:
                    return choice + '.md'

            print(f"Invalid choice. Enter 1-{len(decks)} or a deck name.")

        except KeyboardInterrupt:
            print("\n\nCancelled.")
            sys.exit(0)


def main():
    """Main entry point."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    base_dir = os.path.dirname(script_dir)

    if len(sys.argv) > 1 and not sys.argv[1].startswith('-'):
        # Command line mode
        parser = argparse.ArgumentParser(
            description="Convert Marp markdown to PowerPoint with FASTR styling"
        )
        parser.add_argument('markdown_file', help='Markdown file to convert')
        parser.add_argument('--output', '-o', help='Output PPTX filename')

        args = parser.parse_args()

        md_file = args.markdown_file
        if not md_file.startswith('/'):
            md_file = os.path.join(base_dir, md_file)

        success = convert_to_pptx(md_file, base_dir, args.output)
        sys.exit(0 if success else 1)

    else:
        # Interactive mode
        print("\n" + "=" * 70)
        print("         FASTR POWERPOINT CONVERTER")
        print("=" * 70)

        deck_file = prompt_for_deck(base_dir)
        deck_path = os.path.join(base_dir, "outputs", deck_file)

        success = convert_to_pptx(deck_path, base_dir)
        sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
