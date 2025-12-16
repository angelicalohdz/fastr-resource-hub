#!/usr/bin/env python3
"""
Create FASTR-branded PowerPoint template.
Generates fastr-reference.pptx with proper styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# FASTR brand colors
FASTR_TEAL = RGBColor(15, 112, 109)  # #0f706d
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)

def create_fastr_template():
    """Create a FASTR-branded PowerPoint template."""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Set background to FASTR teal
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = FASTR_TEAL

    # Title
    title = title_slide.shapes.title
    title.text = "Title"
    title_frame = title.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Subtitle
    if len(title_slide.placeholders) > 1:
        subtitle = title_slide.placeholders[1]
        subtitle.text = "Subtitle"
        subtitle_frame = subtitle.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = WHITE

    # Slide 2: Title and Content
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Title
    title = content_slide.shapes.title
    title.text = "Slide Title"
    title_frame = title.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = FASTR_TEAL

    # Body text
    if len(content_slide.placeholders) > 1:
        body = content_slide.placeholders[1]
        body_frame = body.text_frame
        body_frame.text = "Body text"

        # Set font size for body (smaller to prevent overflow)
        for paragraph in body_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = DARK_GRAY
            paragraph.level = 0

    # Save template
    template_path = 'fastr-reference.pptx'
    prs.save(template_path)
    print(f"✅ Created: {template_path}")
    print(f"   FASTR teal: #0f706d")
    print(f"   Title: 44pt (white on teal)")
    print(f"   Headings: 28pt (teal)")
    print(f"   Body: 16pt (dark gray)")
    print(f"\n💡 Use with: python3 convert_to_pptx.py your_deck.md")

if __name__ == "__main__":
    try:
        create_fastr_template()
    except ImportError:
        print("❌ python-pptx not installed")
        print("   Install: pip3 install python-pptx")
